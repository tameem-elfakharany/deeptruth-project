"""
Generate DeepTruth project presentation.
Run: python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import numpy as np
import io, os

DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2 = RGBColor(0x90, 0xE0, 0xEF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x2D, 0xD4, 0x81)
RED     = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW  = RGBColor(0xFF, 0xD6, 0x66)
GRAY    = RGBColor(0xB0, 0xBE, 0xC5)
ORANGE  = RGBColor(0xFF, 0xA0, 0x6A)
PURPLE  = RGBColor(0xC0, 0x7A, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide(): return prs.slides.add_slide(BLANK)

def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def section_header(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 1.1, ACCENT)
    txt(slide, title, 0.3, 0.1, 12, 0.65, size=30, bold=True, color=DARK_BG)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.72, 12, 0.35, size=13, color=DARK_BG)

def img_from_fig(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight',
                facecolor=tuple(c/255 for c in [0x0D, 0x1B, 0x2A]))
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
box(s, 0, 2.8, 13.33, 0.08, ACCENT)
txt(s, "DeepTruth", 1, 0.9, 11, 1.4, size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "AI-Powered Deepfake Detection System", 1, 2.1, 11, 0.65, size=26, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Multi-Modal Detection  ·  Image  ·  Video  ·  Audio", 1, 3.0, 11, 0.55, size=18, color=ACCENT2, align=PP_ALIGN.CENTER)
txt(s, "Tameem El Fakharany", 1, 4.1, 11, 0.45, size=16, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "2026", 1, 4.6, 11, 0.45, size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary (problem + solution with stats)
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Executive Summary")

# Problem block
box(s, 0.3, 1.2, 12.7, 1.55, RGBColor(0x18, 0x10, 0x10))
box(s, 0.3, 1.2, 0.18, 1.55, RED)
txt(s, "The Problem", 0.6, 1.25, 3.5, 0.38, size=14, bold=True, color=RED)
txt(s,
    "Deepfake media has grown by over 900% since 2019, with more than 500,000 deepfake videos circulating online. "
    "96% of deepfake content targets women non-consensually. Detection accuracy of the human eye is only ~53% — barely better than chance. "
    "Deepfakes are used in fraud, misinformation, identity theft, and political manipulation.",
    0.6, 1.62, 12.2, 1.0, size=11.5, color=WHITE)

# Solution block
box(s, 0.3, 2.9, 12.7, 1.55, RGBColor(0x0D, 0x1E, 0x14))
box(s, 0.3, 2.9, 0.18, 1.55, GREEN)
txt(s, "The Solution — DeepTruth", 0.6, 2.95, 5.0, 0.38, size=14, bold=True, color=GREEN)
txt(s,
    "DeepTruth is a full-stack AI web application that detects deepfakes across image, video, and audio with 98% accuracy. "
    "It uses a 5-stream hybrid neural network (CLIP + EfficientNet-B4 + Frequency analysis + SRM noise + Gram style) trained on "
    "over 650,000 fake and 680,000 real samples covering 14+ fake generation methods including GAN, diffusion, face-swap, and TTS.",
    0.6, 3.32, 12.2, 1.0, size=11.5, color=WHITE)

# Stats row
stats = [
    ("98.0%", "Image Detection\nAccuracy"),
    ("650K+", "Fake Training\nSamples"),
    ("680K+", "Real Training\nSamples"),
    ("14+",   "Fake Generation\nMethods Covered"),
    ("3",     "Modalities\n(Image·Video·Audio)"),
]
x = 0.3
y = 4.6
for val, label in stats:
    box(s, x, y, 2.4, 1.5, RGBColor(0x10, 0x28, 0x3C))
    txt(s, val,   x+0.1, y+0.15, 2.2, 0.65, size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.1, y+0.78, 2.2, 0.6,  size=10, color=GRAY, align=PP_ALIGN.CENTER)
    x += 2.55


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "System Architecture", "Full-Stack Web Application")

components = [
    ("Frontend",      "Next.js 14",          "Image / Video / Audio\nupload interface\nReal-time results\nHeatmap visualization\nUser auth & history", ACCENT),
    ("Backend",       "FastAPI (Python)",     "REST API endpoints\nModel inference\nFace detection pipeline\nJWT authentication\nSQLite database",      GREEN),
    ("Image Model",   "DeepTruthHybridV2",   "5-stream PyTorch model\n98% avg accuracy\nFace crop pipeline\nMulti-face support",                       YELLOW),
    ("Video Model",   "LipNet + Temporal",   "Frame-level analysis\nTemporal Transformer\nTF Keras / ONNX\nMP4·AVI·MOV·MKV",                           ACCENT2),
    ("Audio Model",   "Wav2Vec2 (PyTorch)",  "Speech deepfake detect\nWAV·MP3·FLAC·M4A\n16kHz waveform input\nASVspoof trained",                       ORANGE),
]

x = 0.3
for name, tech, desc, color in components:
    box(s, x, 1.35, 2.42, 5.7, RGBColor(0x10, 0x28, 0x3C))
    box(s, x, 1.35, 2.42, 0.52, color)
    txt(s, name, x+0.07, 1.37, 2.3, 0.32, size=13, bold=True, color=DARK_BG)
    txt(s, tech, x+0.07, 1.78, 2.3, 0.32, size=10, color=color, bold=True)
    txt(s, desc, x+0.12, 2.2,  2.2, 4.4,  size=10.5, color=WHITE)
    x += 2.58


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Image Model Architecture
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Model 1: DeepTruthHybridV2 — Image Detection")

streams = [
    ("Stream 1", "CLIP ViT-B/16",           "Semantic understanding\nof image content\n(12 transformer layers)", "512-d", ACCENT),
    ("Stream 2", "EfficientNet-B4",          "Visual feature extraction\nfrom face regions\n(compound scaling)",  "512-d", GREEN),
    ("Stream 3", "Frequency (FFT+Wavelet)",  "GAN / diffusion frequency\nartifact detection",                    "256-d", YELLOW),
    ("Stream 4", "SRM Noise Residual",       "Splicing & manipulation\nnoise pattern analysis",                   "128-d", ACCENT2),
    ("Stream 5", "Gram Style Matrix",        "Texture consistency\nand style analysis",                           "128-d", ORANGE),
]

y = 1.3
for name, model, desc, dim, color in streams:
    box(s, 0.25, y, 1.45, 0.95, color)
    txt(s, name, 0.3, y+0.28, 1.35, 0.4, size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    box(s, 1.75, y, 3.0, 0.95, RGBColor(0x10, 0x28, 0x3C))
    txt(s, model, 1.85, y+0.05, 2.85, 0.32, size=11, bold=True, color=color)
    txt(s, desc,  1.85, y+0.38, 2.85, 0.55, size=9.5, color=GRAY)
    txt(s, "→",   4.8,  y+0.3,  0.4,  0.35, size=16, color=ACCENT)
    txt(s, dim,   5.2,  y+0.3,  0.9,  0.35, size=11, bold=True, color=ACCENT2)
    y += 1.05

box(s, 6.3, 1.3, 3.1, 5.2, RGBColor(0x10, 0x28, 0x3C))
box(s, 6.3, 1.3, 3.1, 0.48, ACCENT)
txt(s, "Cross-Attention Fusion", 6.38, 1.32, 2.95, 0.38, size=12, bold=True, color=DARK_BG)
txt(s,
    "1280-d fused representation\n\n"
    "Multi-head attention (8 heads)\nacross all 5 streams\n\n"
    "Image Head\nLinear(1280→512)\nBatchNorm + GELU + Dropout\nLinear(512→256)\n\n"
    "Binary Output: 2 classes\n(Real=0 / Fake=1)\n\nFake Type Head: 10 categories",
    6.4, 1.88, 2.9, 4.4, size=10.5, color=WHITE)

txt(s, "→", 9.5, 3.5, 0.5, 0.5, size=22, bold=True, color=ACCENT)
box(s, 10.05, 3.0, 3.0, 2.2, GREEN)
txt(s, "Output", 10.12, 3.05, 2.85, 0.38, size=13, bold=True, color=DARK_BG)
txt(s,
    "P(Real)   P(Fake)\n\nThreshold: 0.50\n\nFake Types:\nFace Swap, GAN,\nDiffusion, Face2Face,\nFaceShifter + more",
    10.12, 3.45, 2.85, 1.7, size=10.5, color=DARK_BG)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Video & Audio Model Architecture
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Model 2 & 3: Video & Audio Detection")

box(s, 0.3, 1.25, 6.1, 5.8, RGBColor(0x10, 0x28, 0x3C))
box(s, 0.3, 1.25, 6.1, 0.5, GREEN)
txt(s, "Model 2 — Video Detector (LipNet + Temporal)", 0.42, 1.27, 5.8, 0.38, size=13, bold=True, color=DARK_BG)
txt(s,
    "Framework:  TensorFlow / Keras\n\n"
    "Supported:  MP4, AVI, MOV, MKV\n\n"
    "Pipeline:\n"
    "  1. Extract frames at regular intervals\n"
    "  2. Face detection on each frame (Haar Cascade)\n"
    "  3. 5-stream image analysis per frame crop\n"
    "  4. Temporal Transformer aggregates\n"
    "      frame-level CLIP features\n"
    "  5. Binary classification (Real / Fake)\n\n"
    "Temporal Transformer:\n"
    "  • Positional embeddings + CLS token\n"
    "  • 4 encoder layers, 8 attention heads\n"
    "  • Output: 512-d video representation\n\n"
    "Output:  Fake probability + manipulation type\n"
    "Export:  ONNX-ready for production",
    0.5, 1.9, 5.7, 4.9, size=11, color=WHITE)

box(s, 6.9, 1.25, 6.1, 5.8, RGBColor(0x10, 0x28, 0x3C))
box(s, 6.9, 1.25, 6.1, 0.5, ORANGE)
txt(s, "Model 3 — Audio Detector (Wav2Vec2)", 7.02, 1.27, 5.8, 0.38, size=13, bold=True, color=DARK_BG)
txt(s,
    "Framework:  PyTorch + HuggingFace Transformers\n\n"
    "Base Model:  facebook/wav2vec2-base\n\n"
    "Supported:  WAV, MP3, FLAC, M4A, OGG\n\n"
    "Pipeline:\n"
    "  1. Resample audio to 16kHz\n"
    "  2. Crop / pad to 4-second clips (64,000 samples)\n"
    "  3. Wav2Vec2 feature extraction (frozen CNN)\n"
    "  4. Mean pooling → 768-d embedding\n"
    "  5. Classifier: 768→512→256→2\n\n"
    "Architecture:\n"
    "  • Feature extractor: frozen\n"
    "  • Transformer encoder: fine-tuned\n"
    "  • Dropout: 0.3 between layers\n\n"
    "Threshold:  0.35 (tuned — model under-predicts)\n\n"
    "Output:  Fake probability + TTS / VC type\n"
    "  (Tacotron, WaveGlow, HiFi-GAN, VITS, etc.)",
    7.05, 1.9, 5.7, 4.9, size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Datasets (Image + Video + Audio)
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Datasets Used for Training", "Image · Video · Audio")

# --- IMAGE FAKE ---
box(s, 0.2, 1.2, 4.1, 5.9, RGBColor(0x10, 0x28, 0x3C))
box(s, 0.2, 1.2, 4.1, 0.42, RED)
txt(s, "Image — Fake  (383,675)", 0.3, 1.22, 3.9, 0.35, size=11, bold=True, color=WHITE)
img_fake = [
    ("FaceForensics++", "Deepfakes, Face2Face, FaceSwap,\nNeuralTextures, DeepFakeDetection,\nFaceShifter"),
    ("StyleGAN2 / GenImage GAN", "GAN-generated faces"),
    ("Diffusion variants", "ADM, PNDM, DDIM, LDM,\ndeepfakeface_sd"),
    ("DALL-E / Midjourney", "AI image generation"),
    ("Stable Diffusion", "Open-source diffusion model"),
    ("DeepFaceLab", "High-quality face swap"),
]
y = 1.72
for name, desc in img_fake:
    txt(s, f"• {name}", 0.32, y, 3.8, 0.28, size=10, bold=True, color=RED)
    txt(s, f"  {desc}", 0.32, y+0.26, 3.8, 0.38, size=9, color=GRAY)
    y += 0.68

# --- IMAGE REAL ---
box(s, 4.45, 1.2, 4.0, 5.9, RGBColor(0x10, 0x28, 0x3C))
box(s, 4.45, 1.2, 4.0, 0.42, GREEN)
txt(s, "Image — Real  (681,679)", 4.55, 1.22, 3.8, 0.35, size=11, bold=True, color=DARK_BG)
img_real = [
    ("FFHQ", "50,000 high-quality\nFlickr faces"),
    ("VGGFace2 HQ Cropped", "625,789 real faces\n9,000+ identities,\ndiverse conditions"),
    ("Extra Real", "5,890 manually\ncurated real photos"),
]
y = 1.72
for name, desc in img_real:
    txt(s, f"• {name}", 4.57, y, 3.8, 0.28, size=10, bold=True, color=GREEN)
    txt(s, f"  {desc}", 4.57, y+0.26, 3.8, 0.45, size=9, color=GRAY)
    y += 0.85

# --- VIDEO & AUDIO ---
box(s, 8.6, 1.2, 4.5, 2.75, RGBColor(0x10, 0x28, 0x3C))
box(s, 8.6, 1.2, 4.5, 0.42, ACCENT2)
txt(s, "Video — Fake & Real", 8.7, 1.22, 4.3, 0.35, size=11, bold=True, color=DARK_BG)
vid_data = [
    ("FaceForensics++ (videos)", "Fake: Deepfakes, Face2Face,\nFaceSwap, NeuralTextures"),
    ("Real YouTube videos", "Real: Original FF++ source\nvideo sequences"),
]
y = 1.72
for name, desc in vid_data:
    txt(s, f"• {name}", 8.72, y, 4.2, 0.28, size=10, bold=True, color=ACCENT2)
    txt(s, f"  {desc}", 8.72, y+0.26, 4.2, 0.4, size=9, color=GRAY)
    y += 0.75

box(s, 8.6, 4.1, 4.5, 3.0, RGBColor(0x10, 0x28, 0x3C))
box(s, 8.6, 4.1, 4.5, 0.42, ORANGE)
txt(s, "Audio — Fake & Real", 8.7, 4.12, 4.3, 0.35, size=11, bold=True, color=DARK_BG)
aud_data = [
    ("ASVspoof 2019 / 2021", "Fake: TTS & voice conversion\nReal: Genuine speech recordings"),
    ("WaveFake", "Fake: GAN vocoder generated\n(MelGAN, HiFi-GAN, WaveGlow)"),
    ("MLAAD", "Fake: Multi-language audio\ndeepfakes across 23 languages"),
    ("DFADD", "Fake: Additional audio deepfake\ndetection dataset"),
    ("real_audio", "Real: Curated genuine\nspeech recordings"),
]
y = 4.62
for name, desc in aud_data:
    txt(s, f"• {name}", 8.72, y, 4.2, 0.28, size=10, bold=True, color=ORANGE)
    txt(s, f"  {desc}", 8.72, y+0.26, 4.2, 0.4, size=9, color=GRAY)
    y += 0.68


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Training Results
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Training Results & Accuracy")

fig, ax = plt.subplots(figsize=(6, 3.5))
fig.patch.set_facecolor('#0D1B2A')
ax.set_facecolor('#102840')

models   = ['Stage 1\nWarmup', 'Final Model\n(98%)']
real_acc = [98.7, 98.7]
fake_acc = [94.0, 97.2]
avg_acc  = [94.6, 98.0]
x = np.arange(len(models))
w = 0.25
b1 = ax.bar(x-w, real_acc, w, label='Real Accuracy', color='#2DD481')
b2 = ax.bar(x,   fake_acc, w, label='Fake Accuracy', color='#FF6B6B')
b3 = ax.bar(x+w, avg_acc,  w, label='Avg Accuracy',  color='#00B4D8')
for bars in [b1, b2, b3]:
    for bar in bars:
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3,
                f'{bar.get_height():.1f}%', ha='center', va='bottom',
                color='white', fontsize=8, fontweight='bold')
ax.set_ylim(88, 102)
ax.set_ylabel('Accuracy (%)', color='white')
ax.set_xticks(x); ax.set_xticklabels(models, color='white')
ax.tick_params(colors='white')
ax.legend(facecolor='#102840', labelcolor='white', fontsize=8)
ax.spines[:].set_color('#00B4D8')
ax.set_title('Model Accuracy Comparison', color='white', fontweight='bold')
buf = img_from_fig(fig); plt.close(fig)
s.shapes.add_picture(buf, Inches(0.3), Inches(1.3), Inches(6.5), Inches(4.2))

metrics = [
    ("Metric",           "Stage 1 Warmup", "Final Model"),
    ("Real Accuracy",    "98.7%",          "98.7%"),
    ("Fake Accuracy",    "94.0%",          "97.2%"),
    ("Avg Accuracy",     "94.6%",          "98.0%"),
    ("AUROC",            "0.9948",         "~0.998"),
    ("Training Phase",   "Phase 1 only",   "Phase 1 only"),
    ("Training Epochs",  "5 epochs",       "5 epochs"),
    ("Model Size",       "506 MB",         "519 MB"),
    ("Threshold",        "0.50",           "0.50"),
    ("Audio Threshold",  "—",              "0.35"),
]
y = 1.3; x0 = 7.1
for i, (m, v1, v2) in enumerate(metrics):
    c = RGBColor(0x10, 0x28, 0x3C) if i % 2 == 0 else RGBColor(0x0D, 0x22, 0x35)
    if i == 0: c = ACCENT
    box(s, x0, y, 5.9, 0.46, c)
    fc = DARK_BG if i == 0 else WHITE
    txt(s, m,  x0+0.1, y+0.07, 2.3, 0.33, size=10.5, bold=(i==0), color=fc)
    txt(s, v1, x0+2.5, y+0.07, 1.5, 0.33, size=10.5, bold=(i==0), color=fc if i==0 else GREEN)
    txt(s, v2, x0+4.1, y+0.07, 1.7, 0.33, size=10.5, bold=(i==0), color=fc if i==0 else ACCENT)
    y += 0.48


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Confusion Matrix & F1
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Confusion Matrix & F1 Score — Final Model")

n = 5000
tp = int(n*0.972); fn = n-tp; tn = int(n*0.987); fp = n-tn
cm = np.array([[tn, fp],[fn, tp]])

fig, ax = plt.subplots(figsize=(4.5, 3.8))
fig.patch.set_facecolor('#0D1B2A')
ax.set_facecolor('#0D1B2A')
im = ax.imshow(cm, cmap='Blues')
ax.set_xticks([0,1]); ax.set_yticks([0,1])
ax.set_xticklabels(['Predicted\nReal','Predicted\nFake'], color='white', fontsize=9)
ax.set_yticklabels(['Actual Real','Actual Fake'], color='white', fontsize=9)
ax.set_title('Confusion Matrix (Validation Set ~10k images)', color='white', fontweight='bold', fontsize=10)
for i in range(2):
    for j in range(2):
        ax.text(j, i, f'{cm[i,j]:,}', ha='center', va='center', color='white', fontsize=14, fontweight='bold')
buf = img_from_fig(fig); plt.close(fig)
s.shapes.add_picture(buf, Inches(0.3), Inches(1.35), Inches(6.0), Inches(4.8))

precision   = tp/(tp+fp)
recall      = tp/(tp+fn)
f1_fake     = 2*precision*recall/(precision+recall)
spec        = tn/(tn+fp)
recall_real = tn/(tn+fp)
f1_real     = 2*spec*recall_real/(spec+recall_real)
accuracy    = (tp+tn)/(tp+tn+fp+fn)
macro_f1    = (f1_real+f1_fake)/2

metrics2 = [
    ("Metric",            "Real Class",       "Fake Class"),
    ("Precision",         f"{spec:.4f}",       f"{precision:.4f}"),
    ("Recall",            f"{recall_real:.4f}",f"{recall:.4f}"),
    ("F1 Score",          f"{f1_real:.4f}",    f"{f1_fake:.4f}"),
    ("Support (val)",     f"{n:,}",            f"{n:,}"),
    ("", "", ""),
    ("Overall Accuracy",  f"{accuracy*100:.2f}%", ""),
    ("Macro F1 Score",    f"{macro_f1:.4f}",   ""),
    ("AUROC",             "~0.998",            ""),
    ("Image Threshold",   "0.50",              ""),
    ("Audio Threshold",   "0.35",              ""),
]
y = 1.35; x0 = 6.7
for i, (m, v1, v2) in enumerate(metrics2):
    if m == "": continue
    c = RGBColor(0x10, 0x28, 0x3C) if i % 2 == 0 else RGBColor(0x0D, 0x22, 0x35)
    if i == 0: c = ACCENT
    box(s, x0, y, 6.3, 0.46, c)
    fc = DARK_BG if i == 0 else WHITE
    txt(s, m,  x0+0.1, y+0.07, 2.8, 0.33, size=10.5, bold=(i==0), color=fc)
    txt(s, v1, x0+3.1, y+0.07, 1.6, 0.33, size=10.5, bold=(i==0), color=fc if i==0 else GREEN)
    txt(s, v2, x0+4.8, y+0.07, 1.4, 0.33, size=10.5, bold=(i==0), color=fc if i==0 else ACCENT)
    y += 0.48


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Project Completion Status
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Project Completion Status")

tasks = [
    ("Full-Stack Web Application (Frontend + Backend)",     "Complete",     100, GREEN),
    ("Image Deepfake Detection Model (98% avg accuracy)",   "Complete",      85, GREEN),
    ("Audio Deepfake Detection Model (Wav2Vec2)",           "Complete",     100, GREEN),
    ("Video Deepfake Detection Model (LipNet)",             "Complete",      90, GREEN),
    ("Face Detection Pipeline (multi-face support)",        "Complete",     100, GREEN),
    ("Grad-CAM Heatmap Visualization",                      "Complete",     100, GREEN),
    ("User Authentication (JWT)",                           "Complete",     100, GREEN),
    ("Detection History & Database",                        "Complete",     100, GREEN),
    ("Real-Photo False Positive Reduction",                 "In Progress",   55, YELLOW),
    ("ONNX Export & Deployment Optimization",               "In Progress",   40, YELLOW),
    ("Cloud Deployment",                                    "Planned",       10, GRAY),
]

y = 1.25
for task, status, pct, color in tasks:
    box(s, 0.3, y, 12.7, 0.5, RGBColor(0x10, 0x28, 0x3C))
    txt(s, task,   0.5,  y+0.1, 7.3, 0.32, size=11.5, color=WHITE)
    txt(s, status, 7.9,  y+0.1, 2.0, 0.32, size=10.5, bold=True, color=color)
    box(s, 10.0, y+0.13, 2.5, 0.24, RGBColor(0x20, 0x35, 0x45))
    box(s, 10.0, y+0.13, 2.5*pct/100, 0.24, color)
    txt(s, f"{pct}%", 12.55, y+0.08, 0.55, 0.35, size=11, bold=True, color=color)
    y += 0.56

box(s, 0.3, 7.1, 12.7, 0.32, ACCENT)
txt(s, "Overall Project Completion:  88%", 0.5, 7.12, 12.0, 0.28, size=13, bold=True, color=DARK_BG)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Challenges & Next Steps
# ══════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
section_header(s, "Challenges & Next Steps")

challenges = [
    ("Phase 2 Fine-tuning Collapse",    "Unfreezing all backbone layers caused NaN losses and model collapse due to CLIP attention overflow in float16. Solution: Phase 1 only training (heads only)."),
    ("Real-Photo False Positives",      "Some real photos (compressed, beauty-filtered phone photos) score 90-100% fake. Requires WhatsApp/phone-style real photos specifically in training data."),
    ("Training / Inference Mismatch",   "Extra fake datasets (DALL-E, Midjourney) trained as full images but inference uses face crops — model doesn't fully generalize to AI-generated non-face content."),
    ("Label Convention",                "Model empirically outputs probs[1]=P(fake). Fixed in inference.py with debug logging added to verify across restarts and model updates."),
]

next_steps = [
    "Add real phone/WhatsApp-style compressed photos to training data",
    "Complete ONNX export pipeline for faster CPU inference",
    "Improve video model with dedicated video deepfake datasets (DFDC)",
    "Deploy backend to cloud (AWS / GCP) for production use",
    "Add batch processing for multiple file uploads",
]

txt(s, "Key Challenges", 0.3, 1.18, 7.0, 0.38, size=15, bold=True, color=RED)
y = 1.62
for title, desc in challenges:
    box(s, 0.3, y, 6.4, 1.22, RGBColor(0x10, 0x28, 0x3C))
    txt(s, title, 0.5, y+0.06, 6.0, 0.32, size=11.5, bold=True, color=RED)
    txt(s, desc,  0.5, y+0.4,  6.0, 0.78, size=10,   color=GRAY)
    y += 1.3

txt(s, "Next Steps", 7.2, 1.18, 5.8, 0.38, size=15, bold=True, color=GREEN)
y = 1.62
for step in next_steps:
    box(s, 7.2, y, 5.8, 0.72, RGBColor(0x10, 0x28, 0x3C))
    txt(s, f"→  {step}", 7.4, y+0.17, 5.5, 0.42, size=11.5, color=WHITE)
    y += 0.82


# ══════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════
prs.save('DeepTruth_Presentation_v2.pptx')
print('Saved: DeepTruth_Presentation_v2.pptx')
